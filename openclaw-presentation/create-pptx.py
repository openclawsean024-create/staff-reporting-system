# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(13, 17, 23)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(88, 166, 255)
ORANGE = RGBColor(255, 166, 87)
GRAY = RGBColor(139, 148, 158)

def add_slide(title=None, content=None, is_title=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    if is_title:
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.text = u"OpenClaw 全面介紹"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        sb = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.8))
        p = sb.text_frame.paragraphs[0]
        p.text = u"AI 個人助理的全新可能"
        p.font.size = Pt(22)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    else:
        if title:
            sb = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(4), Inches(0.4))
            p = sb.text_frame.paragraphs[0]
            p.text = u"認識 OpenClaw"
            p.font.size = Pt(10)
            p.font.color.rgb = ORANGE
            p.font.bold = True
            
            tb = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.333), Inches(0.6))
            p = tb.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = WHITE
        
        if content:
            cb = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.333), Inches(5.5))
            tf = cb.text_frame
            tf.word_wrap = True
            for i, line in enumerate(content):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE
                p.space_after = Pt(8)

# 12 Slides
add_slide(is_title=True)
add_slide(u"什麼是 OpenClaw？", [u"核心理念：", u"  - 開源 AI 個人助理框架", u"  - 支援多平台、多訊息源", u"  - 本地部署，隱私優先", u"  - 高度可擴充的技能系統", u"", u"核心能力：", u"  - 跨平台訊息整合", u"  - AI 對話與任務執行", u"  - 自動化工作流程", u"  - Skills Marketplace"])
add_slide(u"六大核心功能", [u"1. 多平台訊息：Telegram, Discord, WhatsApp, Signal", u"2. AI 對話：多模型支援、上下文理解、自訂提示詞", u"3. 自動化：Cron 排程、心跳檢查、工作流程", u"4. 技能系統：Skills Marketplace、自由擴充、社區分享", u"5. 隱私安全：本地部署、端對端加密、自托管", u"6. 開發工具：API、Webhook、CLI"])
add_slide(u"系統架構", [u"應用層：Web / CLI / Telegram / Discord", u"技能層：50+ Skills / MCP Servers", u"核心層：Gateway / Agent / Memory", u"平台層：Telegram / Discord / Slack", u"", u"Gateway：訊息路由、認證授權、平台連接管理", u"Agent：AI 模型調用、對話邏輯、工具調度"])
add_slide(u"跨平台支援", [u"支援平台：", u"  - Telegram, Discord, WhatsApp, Signal", u"  - iMessage, Slack, Email, Web Chat", u"", u"統一入口：透過任何平台都能與你的 AI 助理互動"])
add_slide(u"Skills Marketplace", [u"現有技能分類：", u"  - 圖片生成：OpenAI, Google, MiniMax, DashScope", u"  - 影片生成：Veo, Seedance, Wan", u"  - 搜尋與研究：Web Search, Deep Research", u"  - 生產力：Email, Calendar, Tasks", u"  - 開發工具：GitHub, Docker, Vercel", u"", u"安裝簡單：clawhub install <skill>"])
add_slide(u"多模型支援", [u"支援的 AI 模型提供商：", u"  - OpenAI: GPT-4, GPT-4o, o1, o3", u"  - Google: Gemini 1.5/2.0, Veo, Imagen", u"  - Anthropic: Claude 3.5/3.7, Sonnet, Opus", u"  - MiniMax: M2, M2.1, M3", u"  - 阿里雲：通義千問, 萬象", u"", u"靈活切換：根據任務需求，自動或手動切換最適合的 AI 模型"])
add_slide(u"應用場景", [u"個人助理：郵件管理、日曆安排、任務提醒、待辦事項", u"創意工具：AI 繪圖、影片生成、簡報製作、內容創作", u"開發者：程式碼管理、部署上線、API 整合、自動化測試", u"研究與分析：網路搜尋、資料分析、報告生成、趨勢追蹤", u"監控與警報：系統監控、安全警報、定時任務、異常通知"])
add_slide(u"快速開始", [u"安裝步驟：", u"  1. npm install -g openclaw", u"  2. openclaw init", u"  3. 設定 Token 連接平台", u"  4. openclaw start", u"", u"系統需求：Node.js 18+, 4GB+ RAM, 支援 Docker", u"部署選項：本地部署、Docker Compose、雲端 VPS"])
add_slide(u"社群與生態", [u"資源與支援：", u"  - 官方文檔：docs.openclaw.ai", u"  - Discord 社群", u"  - GitHub Issues", u"  - ClawHub 技能市場", u"", u"參與貢獻：", u"  - 提交技能到 ClawHub", u"  - 回報問題與功能請求", u"  - 翻譯與文檔"])
add_slide(u"未來規劃", [u"短期目標 (0-3個月)：", u"  - 完善技能市場生態", u"  - 增加更多 AI 模型支援", u"  - 改進 MCP 整合", u"", u"中期目標 (3-6個月)：", u"  - 多語言支援優化", u"  - 行動應用程式", u"  - 進階自動化工作流", u"", u"長期願景 (6-12個月)：", u"  - 完全自主的 AI 代理系統", u"  - 跨裝置無縫同步"])
add_slide(u"結語", [u"「OpenClaw — 讓 AI 助理真正成為你的數位延伸，", u"跨越所有平台，服務於你的日常生活與工作效率。」", u"", u"— OpenClaw 團隊", u"", u"立即開始：", u"  - GitHub: github.com/openclaw/openclaw", u"  - Discord: discord.com/invite/clawd", u"  - 文檔: docs.openclaw.ai"])

prs.save(r"C:\Users\sean\.openclaw\workspace\openclaw-presentation\OpenClaw_Presentation_CN.pptx")
print("Done!")
